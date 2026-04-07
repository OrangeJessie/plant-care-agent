"""growth_slides — 将植物成长日志生成 PPTX 幻灯片。

从 data/garden/{plant}.md 读取日志和 data/charts/ 下的图表，
自动组装成一份美观的成长故事 PPT。

幻灯片结构:
  1. 封面：植物名 + 种植天数 + 阶段
  2. 基本信息：品种/位置/种植日期等
  3-N. 按时间分组的日志页（每页一天或一周）
  N+1. 图表页（如果有 charts 目录下的图）
  N+2. 总结页
"""

import logging
import re
from collections import defaultdict
from datetime import datetime
from pathlib import Path

from pydantic import Field

from nat.builder.builder import Builder
from nat.builder.function_info import FunctionInfo
from nat.cli.register_workflow import register_function
from nat.data_models.function import FunctionBaseConfig

logger = logging.getLogger(__name__)

FRONTMATTER_RE = re.compile(r"^---\s*\n(.*?)\n---\s*\n", re.DOTALL)
EVENT_LINE_RE = re.compile(r"^- \*\*\[(.+?)]\*\*\s+(.+?)(?:\s+_(\d{2}:\d{2})_)?$", re.MULTILINE)
DATE_HEADING_RE = re.compile(r"^### (\d{4}-\d{2}-\d{2})", re.MULTILINE)

STAGE_EMOJI = {
    "播种期": "🌰", "苗期": "🌱", "生长期": "🌿",
    "花期": "🌸", "结果期": "🍅", "采收期": "🎉",
}


class GrowthSlidesConfig(FunctionBaseConfig, name="growth_slides"):
    garden_dir: str = Field(default="./data/garden")
    charts_dir: str = Field(default="./data/charts")
    output_dir: str = Field(default="./data/slides")


def _parse_fm(text: str) -> dict[str, str]:
    m = FRONTMATTER_RE.match(text)
    if not m:
        return {}
    out: dict[str, str] = {}
    for line in m.group(1).splitlines():
        if ":" in line:
            k, v = line.split(":", 1)
            out[k.strip()] = v.strip()
    return out


def _parse_events(body: str) -> list[dict]:
    events: list[dict] = []
    current_date = ""
    for line in body.splitlines():
        dm = DATE_HEADING_RE.match(line)
        if dm:
            current_date = dm.group(1)
            continue
        em = EVENT_LINE_RE.match(line)
        if em and current_date:
            events.append({
                "date": current_date,
                "type": em.group(1),
                "desc": em.group(2).strip(),
                "time": em.group(3) or "",
            })
    return events


@register_function(config_type=GrowthSlidesConfig)
async def growth_slides_function(config: GrowthSlidesConfig, _builder: Builder):
    garden = Path(config.garden_dir)
    charts_dir = Path(config.charts_dir)
    output_dir = Path(config.output_dir)

    async def _generate_slides(plant_id: str) -> str:
        """Generate a PPTX slide deck showing a plant's growth story.
        Input: plant name/ID.
        Returns: path to the generated .pptx file.
        The slides include cover, info, daily logs, charts (if available), and summary."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        except ImportError:
            return "需要安装 python-pptx: pip install python-pptx"

        pid = plant_id.strip()
        path = garden / f"{pid}.md"
        if not path.exists():
            return f"未找到「{pid}」的日志文件。"

        text = path.read_text(encoding="utf-8")
        fm = _parse_fm(text)
        body = FRONTMATTER_RE.sub("", text, count=1)
        events = _parse_events(body)

        if not events:
            return f"「{pid}」暂无事件记录，无法生成幻灯片。"

        name = fm.get("name", pid)
        species = fm.get("species", "-")
        location = fm.get("location", "-")
        planted = fm.get("planted", "-")
        stage = fm.get("stage", "-")
        stage_emoji = STAGE_EMOJI.get(stage, "🌱")

        try:
            days = (datetime.now().date() - datetime.strptime(planted, "%Y-%m-%d").date()).days
        except (ValueError, TypeError):
            days = "?"

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        GREEN_DARK = RGBColor(0x2E, 0x7D, 0x32)
        GREEN_LIGHT = RGBColor(0x81, 0xC7, 0x84)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        GRAY = RGBColor(0x61, 0x61, 0x61)
        BG_COLOR = RGBColor(0xFA, 0xFA, 0xF5)

        def _set_slide_bg(slide, color=BG_COLOR):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def _add_text(slide, left, top, width, height, text_val, font_size=18,
                      bold=False, color=GRAY, alignment=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                             Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text_val
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = alignment
            return tf

        # --- Slide 1: Cover ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _set_slide_bg(slide, RGBColor(0x2E, 0x7D, 0x32))
        _add_text(slide, 1, 1.5, 11, 1.5,
                  f"{stage_emoji} {name}", 44, True, WHITE, PP_ALIGN.CENTER)
        _add_text(slide, 1, 3.2, 11, 1,
                  f"品种: {species}  |  位置: {location}", 20, False, GREEN_LIGHT, PP_ALIGN.CENTER)
        _add_text(slide, 1, 4.5, 11, 1,
                  f"种植天数: {days}天  |  当前阶段: {stage}  |  记录: {len(events)}条",
                  18, False, GREEN_LIGHT, PP_ALIGN.CENTER)
        _add_text(slide, 1, 6, 11, 0.5,
                  f"生成于 {datetime.now().strftime('%Y-%m-%d %H:%M')}",
                  12, False, GREEN_LIGHT, PP_ALIGN.CENTER)

        # --- Slide 2: Info ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide)
        _add_text(slide, 0.8, 0.5, 11, 0.8, "📋 基本信息", 28, True, GREEN_DARK)

        info_lines = [
            f"🌿 植物名称:  {name}",
            f"🧬 品种:  {species}",
            f"📍 种植位置:  {location}",
            f"📅 种植日期:  {planted}",
            f"⏱️ 种植天数:  {days}天",
            f"📊 当前阶段:  {stage}",
            f"📝 总事件数:  {len(events)}条",
        ]
        _add_text(slide, 1.2, 1.8, 10, 4.5, "\n\n".join(info_lines), 18, False, GRAY)

        # --- Slides 3+: Daily logs grouped by date ---
        by_date: dict[str, list[dict]] = defaultdict(list)
        for ev in events:
            by_date[ev["date"]].append(ev)

        date_keys = sorted(by_date.keys())
        page_size = 5
        for page_start in range(0, len(date_keys), page_size):
            page_dates = date_keys[page_start:page_start + page_size]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_bg(slide)

            range_label = page_dates[0]
            if len(page_dates) > 1:
                range_label = f"{page_dates[0]} ~ {page_dates[-1]}"
            _add_text(slide, 0.8, 0.4, 11, 0.7, f"📖 {range_label}", 24, True, GREEN_DARK)

            y = 1.4
            for d in page_dates:
                _add_text(slide, 1, y, 10, 0.4, f"📅 {d}", 16, True, GREEN_DARK)
                y += 0.5
                for ev in by_date[d]:
                    time_str = f" ({ev['time']})" if ev['time'] else ""
                    line = f"  • [{ev['type']}] {ev['desc']}{time_str}"
                    _add_text(slide, 1.3, y, 10, 0.35, line, 14, False, GRAY)
                    y += 0.4
                y += 0.2

        # --- Chart slides ---
        chart_files = []
        if charts_dir.is_dir():
            for suffix in ("_timeline.png", "_dashboard.png"):
                cf = charts_dir / f"{pid}{suffix}"
                if cf.exists():
                    chart_files.append(cf)

        for cf in chart_files:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_bg(slide)
            chart_label = "时间线" if "timeline" in cf.name else "综合看板"
            _add_text(slide, 0.8, 0.3, 11, 0.6, f"📊 {chart_label}", 24, True, GREEN_DARK)
            slide.shapes.add_picture(str(cf), Inches(0.8), Inches(1.2), Inches(11.5))

        # --- Summary slide ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide, GREEN_DARK)
        _add_text(slide, 1, 2, 11, 1.5,
                  f"🌟 {name} 的成长旅程", 36, True, WHITE, PP_ALIGN.CENTER)

        from collections import Counter
        type_counts = Counter(e["type"] for e in events)
        top3 = type_counts.most_common(3)
        summary_parts = [f"{t}: {c}次" for t, c in top3]
        _add_text(slide, 1, 4, 11, 1,
                  f"共记录 {len(events)} 条养护事件，跨越 {len(date_keys)} 天",
                  20, False, GREEN_LIGHT, PP_ALIGN.CENTER)
        _add_text(slide, 1, 5, 11, 0.8,
                  f"主要活动: {' | '.join(summary_parts)}",
                  18, False, GREEN_LIGHT, PP_ALIGN.CENTER)
        _add_text(slide, 1, 6.2, 11, 0.5,
                  "继续加油，见证每一天的成长 🌻",
                  16, False, GREEN_LIGHT, PP_ALIGN.CENTER)

        # --- Save ---
        output_dir.mkdir(parents=True, exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d_%H%M")
        out_path = output_dir / f"{pid}_成长故事_{ts}.pptx"
        prs.save(str(out_path))

        return (
            f"✅ 幻灯片已生成: {out_path}\n"
            f"   共 {len(prs.slides)} 页，包含封面、信息、"
            f"{len(date_keys)} 天日志"
            f"{'、' + str(len(chart_files)) + '张图表' if chart_files else ''}"
            f"和总结页。"
        )

    async def _generate_garden_slides(query: str) -> str:
        """Generate a PPTX slide deck for ALL plants in the garden.
        Input: anything (ignored). Creates a multi-plant overview deck.
        Returns: path to the generated .pptx file."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return "需要安装 python-pptx: pip install python-pptx"

        files = sorted(p for p in garden.glob("*.md") if p.name != "GARDEN.md")
        if not files:
            return "暂无植物日志，无法生成幻灯片。"

        GREEN_DARK = RGBColor(0x2E, 0x7D, 0x32)
        GREEN_LIGHT = RGBColor(0x81, 0xC7, 0x84)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        GRAY = RGBColor(0x61, 0x61, 0x61)
        BG_COLOR = RGBColor(0xFA, 0xFA, 0xF5)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        def _set_bg(slide, color=BG_COLOR):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def _txt(slide, left, top, width, height, text_val,
                 font_size=18, bold=False, color=GRAY, align=PP_ALIGN.LEFT):
            box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                           Inches(width), Inches(height))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text_val
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = align

        # Cover
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, GREEN_DARK)
        _txt(slide, 1, 2, 11, 1.5, "🌱 我的花园", 44, True, WHITE, PP_ALIGN.CENTER)
        _txt(slide, 1, 4, 11, 1, f"共 {len(files)} 株植物的成长记录",
             22, False, GREEN_LIGHT, PP_ALIGN.CENTER)
        _txt(slide, 1, 5.5, 11, 0.5,
             f"生成于 {datetime.now().strftime('%Y-%m-%d %H:%M')}",
             12, False, GREEN_LIGHT, PP_ALIGN.CENTER)

        # Per-plant summary
        for fpath in files:
            text = fpath.read_text(encoding="utf-8")
            fm = _parse_fm(text)
            body = FRONTMATTER_RE.sub("", text, count=1)
            events = _parse_events(body)

            pname = fm.get("name", fpath.stem)
            stage = fm.get("stage", "-")
            emoji = STAGE_EMOJI.get(stage, "🌱")

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(slide)
            _txt(slide, 0.8, 0.5, 11, 0.8, f"{emoji} {pname}", 28, True, GREEN_DARK)

            info = (
                f"品种: {fm.get('species', '-')}  |  "
                f"位置: {fm.get('location', '-')}  |  "
                f"阶段: {stage}  |  "
                f"记录: {len(events)}条"
            )
            _txt(slide, 1, 1.5, 11, 0.5, info, 16, False, GRAY)

            if events:
                recent = events[-5:]
                y = 2.5
                _txt(slide, 1, y, 10, 0.5, "最近记录:", 16, True, GREEN_DARK)
                y += 0.6
                for ev in recent:
                    line = f"  {ev['date']} [{ev['type']}] {ev['desc']}"
                    _txt(slide, 1.2, y, 10, 0.35, line, 14, False, GRAY)
                    y += 0.45

            chart_path = charts_dir / f"{fpath.stem}_dashboard.png"
            if chart_path.exists():
                try:
                    slide.shapes.add_picture(str(chart_path), Inches(0.8), Inches(4.5), Inches(11.5))
                except Exception:
                    pass

        output_dir.mkdir(parents=True, exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d_%H%M")
        out_path = output_dir / f"花园总览_{ts}.pptx"
        prs.save(str(out_path))

        return f"✅ 花园总览幻灯片已生成: {out_path}\n   共 {len(prs.slides)} 页。"

    yield FunctionInfo.from_fn(
        _generate_slides,
        description="为单株植物生成成长故事 PPTX 幻灯片（封面+信息+日志+图表+总结）。",
    )
    yield FunctionInfo.from_fn(
        _generate_garden_slides,
        description="为花园所有植物生成总览 PPTX 幻灯片。",
    )
